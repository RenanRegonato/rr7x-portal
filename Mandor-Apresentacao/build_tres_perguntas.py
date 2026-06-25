# -*- coding: utf-8 -*-
"""
Gera 'Mandor-Tres-Perguntas.pptx' a partir da narrativa de 6 slides
(conversa entre pares institucionais, não pitch). Identidade editorial 2026
do Mandor: Archive Ivory / Editorial Ink / Oxidized Bronze, serifa Newsreader.
Regras: sem travessao, "O Mandor" (masculino), "rede cognitiva" (nunca software).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Paleta ----
INK        = RGBColor(0x18, 0x14, 0x0F)   # espresso quente (fundo escuro)
INK_WARM   = RGBColor(0x22, 0x1B, 0x13)   # ponto mais quente do gradiente
IVORY      = RGBColor(0xED, 0xE9, 0xE5)   # Archive Ivory
IVORY_DIM  = RGBColor(0xB7, 0xB0, 0xA6)   # marfim atenuado (apoio sobre escuro)
BRONZE     = RGBColor(0x8C, 0x6F, 0x45)   # Oxidized Bronze
BRONZE_LT  = RGBColor(0xC0, 0xA2, 0x6C)   # bronze claro (acento sobre escuro)
INK_TEXT   = RGBColor(0x21, 0x1F, 0x1A)   # Editorial Ink (texto sobre claro)
BODY_GRAY  = RGBColor(0x5E, 0x5A, 0x53)   # corpo sobre claro
LIGHT_BG   = RGBColor(0xED, 0xE9, 0xE5)
LIGHT_DEEP = RGBColor(0xE2, 0xDB, 0xD1)   # painel/marfim mais profundo

SERIF = "Newsreader"
SANS  = "Hanken Grotesk"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------

def add_slide(dark=True):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.gradient()
    stops = fill.gradient_stops
    if dark:
        stops[0].color.rgb = INK_WARM; stops[0].position = 0.0
        stops[1].color.rgb = INK;      stops[1].position = 1.0
    else:
        stops[0].color.rgb = LIGHT_BG; stops[0].position = 0.0
        stops[1].color.rgb = LIGHT_DEEP; stops[1].position = 1.0
    try:
        fill.gradient_angle = 55.0
    except Exception:
        pass
    return s


def _set_spc(run, spc):
    if spc:
        run._r.get_or_add_rPr().set('spc', str(int(spc)))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    return tb, tf


def style_para(p, runs, align=PP_ALIGN.LEFT, space_after=0, space_before=0, line=None):
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    for text, size, color, font, bold, italic, spc in runs:
        r = p.add_run(); r.text = text
        f = r.font
        f.size = Pt(size); f.name = font; f.bold = bold; f.italic = italic
        f.color.rgb = color
        _set_spc(r, spc)
    return p


def label(slide, x, y, text, color=BRONZE, rule=True):
    """Rotulo caixa-alta com tracking e tracinho bronze, como no deck atual."""
    if rule:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Pt(5), Inches(0.32), Pt(2))
        ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background()
        tx = x + Inches(0.46)
    else:
        tx = x
    tb, tf = textbox(slide, tx, y, Inches(9.5), Inches(0.35))
    p = tf.paragraphs[0]
    style_para(p, [(text, 11.5, color, SANS, True, False, 320)])
    return tb


def monogram(slide, ox, oy, h_in, color):
    """Monograma I/I (viewBox 48x50) recriado em vetor nativo."""
    s = Inches(h_in) / 50.0  # EMU por unidade svg

    def rect(x, w):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    int(ox + x * s), int(oy),
                                    int(w * s), int(50 * s))
        sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
        sp.shadow.inherit = False
        return sp
    rect(0, 10)
    rect(37.4, 10)
    # barra inclinada (poligono)
    pts = [(25.4, 0), (35.4, 0), (22, 50), (12, 50)]
    start = (int(ox + pts[0][0] * s), int(oy + pts[0][1] * s))
    fb = slide.shapes.build_freeform(start[0], start[1], scale=1.0)
    fb.add_line_segments([(int(ox + px * s), int(oy + py * s)) for px, py in pts[1:]],
                         close=True)
    sh = fb.convert_to_shape()
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False


def wordmark(slide, x, y, w, color, size=22, align=PP_ALIGN.LEFT):
    tb, tf = textbox(slide, x, y, w, Inches(0.6))
    style_para(tf.paragraphs[0], [("MANDOR", size, color, SERIF, False, False, 620)], align=align)
    return tb


def footer(slide, dark=True, page=None, total=6):
    color = IVORY_DIM if dark else BODY_GRAY
    mono = IVORY if dark else INK_TEXT
    y = SH - Inches(0.62)
    monogram(slide, Inches(0.7), y, 0.20, mono)
    tb, tf = textbox(slide, Inches(0.95), y - Pt(2), Inches(5), Inches(0.3))
    style_para(tf.paragraphs[0], [("MANDOR", 10, color, SANS, True, False, 300)])
    if page is not None:
        tb2, tf2 = textbox(slide, SW - Inches(2.2), y - Pt(2), Inches(1.5), Inches(0.3))
        style_para(tf2.paragraphs[0], [(f"{page:02d} / {total:02d}", 10, color, SANS, False, False, 200)],
                   align=PP_ALIGN.RIGHT)


def notes(slide, apoio, psico, visual):
    tf = slide.notes_slide.notes_text_frame
    tf.text = "TEXTO DE APOIO PARA O APRESENTADOR\n" + apoio
    for titulo, corpo in [("OBJETIVO PSICOLÓGICO DO SLIDE", psico),
                          ("SUGESTÃO VISUAL", visual)]:
        p = tf.add_paragraph(); p.text = ""
        p = tf.add_paragraph(); p.text = titulo
        p.font.bold = True
        p = tf.add_paragraph(); p.text = corpo


MX = Inches(0.95)          # margem esquerda padrão
CW = SW - Inches(1.9)      # largura util

# =================================================================
# SLIDE 1 - Abertura: as tres perguntas
# =================================================================
s = add_slide(dark=True)
label(s, MX, Inches(0.85), "PARA COMEÇAR", BRONZE_LT)
tb, tf = textbox(s, MX, Inches(1.45), Inches(10.6), Inches(1.9))
style_para(tf.paragraphs[0],
           [("Três perguntas que surgem em qualquer ", 41, IVORY, SERIF, False, False, 0),
            ("operação que cresce", 41, BRONZE_LT, SERIF, False, True, 0)],
           line=1.04)
tb, tf = textbox(s, MX, Inches(3.25), Inches(9.6), Inches(1.4))
style_para(tf.paragraphs[0],
           [("Toda instituição financeira que amadurece chega, mais cedo ou mais tarde, ao mesmo "
             "conjunto de perguntas. Não são sinais de fragilidade. São as perguntas que distinguem "
             "quem cresce de quem cresce com método.", 15.5, IVORY_DIM, SERIF, False, False, 0)],
           line=1.28)
# preview das tres perguntas
qy = Inches(4.75)
perguntas = [
    "Quando um profissional-chave sai, o método permanece?",
    "O mercado está mapeado ou ainda depende de relacionamento?",
    "A capacidade de análise acompanha o volume de oportunidades?",
]
for i, q in enumerate(perguntas):
    tb, tf = textbox(s, MX, qy + Inches(0.62) * i, Inches(11.2), Inches(0.6))
    p = tf.paragraphs[0]
    style_para(p, [(f"0{i+1}", 17, BRONZE_LT, SERIF, False, False, 0),
                   ("    " + q, 16, IVORY, SERIF, False, False, 0)], line=1.0)
footer(s, dark=True, page=1)
notes(s,
      "Abra sem slides na cabeça. Diga que não veio apresentar nada, veio propor uma conversa. "
      "Ao longo dos anos, conversando com casas em estagios diferentes, você percebeu que tres "
      "perguntas aparecem em praticamente toda instituição que cresce. Nao porque algo esteja errado, "
      "mas porque crescer expoe a estrutura. Convide-os a responder mentalmente, não em voz alta.",
      "Estabelecer paridade entre pares e desarmar a expectativa de pitch. Posicionar você como quem "
      "observa o mercado, não como quem vende. Criar a moldura: estas perguntas são legitimas e "
      "universais, logo reconhece-las não e admitir falha.",
      "Fundo espresso, respiro amplo, as tres perguntas em numerais bronze. Sem ícones, sem gráfico.")

# =================================================================
# SLIDES 2-4 - As tres perguntas, uma por slide
# =================================================================
QSLIDES = [
    dict(
        n="PERGUNTA 01", tema="A pergunta da memória institucional.",
        titulo="Quando um profissional-chave sai, o método permanece?",
        subs=["O critério de análise vive nas pessoas ou na instituição?",
              "O conhecimento acumulado é transferível, ou recomeça a cada saída?",
              "A senioridade é um ativo da casa, ou um risco concentrado em poucos nomes?"],
        apoio="Toda casa madura tem nomes que carregam o critério. Isso e força, e também concentração. "
              "Pergunte, sem cobrar: o que sai junto com a pessoa? O método esta documentado, ou esta na "
              "cabeça de quem decide há vinte anos? Nao ha resposta certa, ha graus. Deixe o silêncio trabalhar.",
        psico="Ativar a consciência de risco de concentração sem acusar. Fazer o sócio sênior pensar no "
              "próprio sucessor. Plantar a ideia de que método e patrimônio da instituição, não rotina.",
        visual="Pergunta em marfim ocupando o terço superior. Três subperguntas discretas em bronze. "
               "Nenhuma resposta na tela: o vazio e proposital.",
    ),
    dict(
        n="PERGUNTA 02", tema="A pergunta da originação.",
        titulo="O mercado está mapeado ou ainda depende de relacionamento?",
        subs=["A originação escala com a casa, ou com a agenda de cada sócio?",
              "A cada operação, sabemos exatamente para quem levá-la?",
              "O relacionamento é o diferencial da casa, ou o seu teto?"],
        apoio="Relacionamento e o ativo mais nobre desse mercado, e o mais dificil de escalar. Pergunte: "
              "quanto da originação depende de quem você conhece pessoalmente? Quando um deal chega, o mapa "
              "de para quem levá-lo esta na cabeça de alguém ou na instituição? Reconheca o valor do "
              "relacionamento antes de questionar o seu limite.",
        psico="Mostrar que relacionamento e, ao mesmo tempo, vantagem e limite de escala, sem nomear o limite. "
              "Levar a pensar em originação e cobertura de mercado como capacidade institucional, não pessoal.",
        visual="Mesma gramatica visual da pergunta anterior, para criar ritmo de tríade.",
    ),
    dict(
        n="PERGUNTA 03", tema="A pergunta da escala.",
        titulo="A capacidade de análise acompanha o volume de oportunidades?",
        subs=["Cada nova oportunidade exige mais gente, ou o mesmo critério?",
              "O padrão de análise se sustenta quando o fluxo dobra?",
              "Crescer em volume amplia a consistência, ou a dilui?"],
        apoio="O crescimento de fluxo costuma chegar antes do crescimento de equipe. Pergunte: quando o "
              "volume dobra, o que acontece com a profundidade da análise? Contrata-se, ou aceita-se filtrar "
              "pior? O ponto não e fazer mais rapido, e manter o mesmo critério institucional em escala.",
        psico="Conectar volume a consistência, não a velocidade. Evitar a leitura simplista de eficiência. "
              "Posicionar a capacidade analítica como atributo da instituição que precisa crescer sem se diluir.",
        visual="Fecha a tríade no mesmo peso visual das duas anteriores.",
    ),
]
for i, q in enumerate(QSLIDES):
    s = add_slide(dark=True)
    label(s, MX, Inches(0.9), q["n"], BRONZE_LT)
    # tema curto acima da pergunta
    tb, tf = textbox(s, MX, Inches(1.5), Inches(11), Inches(0.45))
    style_para(tf.paragraphs[0], [(q["tema"], 15, BRONZE_LT, SERIF, False, True, 0)])
    # pergunta gigante
    tb, tf = textbox(s, MX, Inches(2.1), Inches(11.2), Inches(2.2), anchor=MSO_ANCHOR.TOP)
    style_para(tf.paragraphs[0], [(q["titulo"], 44, IVORY, SERIF, False, False, 0)], line=1.05)
    # subperguntas
    sy = Inches(4.65)
    for j, sub in enumerate(q["subs"]):
        tb, tf = textbox(s, MX, sy + Inches(0.6) * j, Inches(11), Inches(0.55))
        p = tf.paragraphs[0]
        style_para(p, [("/  ", 15, BRONZE, SANS, True, False, 0),
                       (sub, 16, IVORY_DIM, SERIF, False, False, 0)], line=1.0)
    footer(s, dark=True, page=2 + i)
    notes(s, q["apoio"], q["psico"], q["visual"])

# =================================================================
# SLIDE 5 - O estado resolvido (fundo claro)
# =================================================================
s = add_slide(dark=False)
label(s, MX, Inches(0.85), "O ESTADO RESOLVIDO", BRONZE)
tb, tf = textbox(s, MX, Inches(1.45), Inches(11.3), Inches(1.6))
style_para(tf.paragraphs[0],
           [("O que acontece quando essas três perguntas estão ", 36, INK_TEXT, SERIF, False, False, 0),
            ("resolvidas?", 36, BRONZE, SERIF, False, True, 0)], line=1.06)
outcomes = [
    "Inteligência institucionalizada, que não depende de quem está na sala.",
    "Mercado mapeado, com destino claro para cada operação.",
    "Capacidade de análise que escala sem inflar a equipe.",
    "Decisões mais consistentes e menos dependentes do indivíduo.",
    "Mais capacidade de execução, com o critério da casa preservado.",
]
oy = Inches(3.3)
for i, o in enumerate(outcomes):
    # tick bronze
    tk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, oy + Inches(0.66) * i + Pt(4), Inches(0.22), Pt(2.4))
    tk.fill.solid(); tk.fill.fore_color.rgb = BRONZE; tk.line.fill.background(); tk.shadow.inherit = False
    tb, tf = textbox(s, MX + Inches(0.42), oy + Inches(0.66) * i, Inches(10.6), Inches(0.6))
    style_para(tf.paragraphs[0], [(o, 17.5, INK_TEXT, SERIF, False, False, 0)], line=1.05)
footer(s, dark=False, page=5)
notes(s,
      "Mude o tom: do questionamento para a imagem do estado desejado. Nao prometa, descreva. Diga que as "
      "casas que respondem bem a essas tres perguntas tendem a compartilhar esses cinco traços. Deixe que "
      "eles se projetem nessa imagem. Ainda não mencione o Mandor.",
      "Criar desejo pelo estado final antes de apresentar qualquer resposta. Transformar perguntas em visão, "
      "de modo que a solucao, quando vier, pareca óbvia e não vendida.",
      "Vira para fundo claro (Archive Ivory). A mudanca de fundo sinaliza resolucao e clareza. Marcadores em bronze.")

# =================================================================
# SLIDE 6 - Por que começamos a construir o Mandor (fundo escuro)
# =================================================================
s = add_slide(dark=True)
label(s, MX, Inches(0.85), "POR QUE EXISTIMOS", BRONZE_LT)
# monograma + wordmark
monogram(s, MX, Inches(1.55), 0.62, IVORY)
wordmark(s, MX + Inches(0.78), Inches(1.55), Inches(6), IVORY, size=30)
tb, tf = textbox(s, MX, Inches(2.7), Inches(11), Inches(1.0))
style_para(tf.paragraphs[0],
           [("Por que começamos a construir ", 34, IVORY, SERIF, False, False, 0),
            ("o Mandor", 34, BRONZE_LT, SERIF, False, True, 0)], line=1.05)
tb, tf = textbox(s, MX, Inches(3.85), Inches(10.7), Inches(2.3))
p = tf.paragraphs[0]
style_para(p, [("Vimos esse padrão se repetir, instituição após instituição. O conhecimento que não fica "
                "quando alguém sai. O mercado que ainda depende de quem se conhece. A análise que não "
                "acompanha o volume.", 16.5, IVORY_DIM, SERIF, False, False, 0)], line=1.3)
p2 = tf.add_paragraph()
style_para(p2, [("O Mandor", 16.5, IVORY, SERIF, False, False, 0),
                (" nasceu dessa observação. Uma rede cognitiva para que método, mercado e capacidade de "
                 "análise deixem de viver em indivíduos e passem a pertencer à instituição.",
                 16.5, IVORY, SERIF, False, False, 0)], line=1.3, space_before=12)
# fechamento
tb, tf = textbox(s, MX, Inches(6.3), Inches(8), Inches(0.5))
style_para(tf.paragraphs[0], [("www.mandor.com.br", 16, BRONZE_LT, SANS, True, False, 120)])
footer(s, dark=True, page=6)
notes(s,
      "So agora o nome aparece. Nao mostre telas, módulos ou funcionalidades. Diga apenas que o Mandor "
      "nasceu de observar esse padrão repetidamente e que foi construido para responder exatamente a essas "
      "tres perguntas. Convide para uma conversa de continuidade, não para uma demonstração. Encerre "
      "devolvendo a palavra a eles.",
      "Apresentar o Mandor como consequência natural da conversa, e não como produto. Ancorar a marca nas "
      "tres perguntas que eles acabaram de reconhecer como suas. Posiciona-lo como rede cognitiva entre "
      "pares, nunca como software a ser vendido.",
      "Retorna ao espresso. Monograma I/I e wordmark MANDOR no topo. Texto sóbrio, fechamento com o "
      "endereço www.mandor.com.br.")

out = "/Users/renan/Desktop/rr7x-portal/Mandor-Apresentacao/Mandor-Tres-Perguntas.pptx"
prs.save(out)
print("OK ->", out)
