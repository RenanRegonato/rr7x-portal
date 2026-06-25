# -*- coding: utf-8 -*-
"""
Deck enxuto para a reuniao com a Assis Gestao (M&A) - Alexandre Assis + Jamerson Marra.
Backbone de conversa (nao pitch): fit honesto 5-de-12, as 3 perguntas, demo ao vivo
no Lumina, estado resolvido, piloto primeiro (sem numeros). Identidade editorial 2026.
Anchors do apresentador vivem nas NOTAS de cada slide (fala com as proprias palavras).
Regras: sem travessao, "O Mandor" (masculino), "rede cognitiva" (nunca software).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK        = RGBColor(0x18, 0x14, 0x0F)
INK_WARM   = RGBColor(0x22, 0x1B, 0x13)
IVORY      = RGBColor(0xED, 0xE9, 0xE5)
IVORY_DIM  = RGBColor(0xB7, 0xB0, 0xA6)
BRONZE     = RGBColor(0x8C, 0x6F, 0x45)
BRONZE_LT  = RGBColor(0xC0, 0xA2, 0x6C)
INK_TEXT   = RGBColor(0x21, 0x1F, 0x1A)
BODY_GRAY  = RGBColor(0x5E, 0x5A, 0x53)
LIGHT_BG   = RGBColor(0xED, 0xE9, 0xE5)
LIGHT_DEEP = RGBColor(0xE2, 0xDB, 0xD1)
SERIF = "Newsreader"
SANS  = "Hanken Grotesk"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
TOTAL = 7

def add_slide(dark=True):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill; fill.gradient()
    st = fill.gradient_stops
    if dark:
        st[0].color.rgb = INK_WARM; st[0].position = 0.0
        st[1].color.rgb = INK;      st[1].position = 1.0
    else:
        st[0].color.rgb = LIGHT_BG; st[0].position = 0.0
        st[1].color.rgb = LIGHT_DEEP; st[1].position = 1.0
    try: fill.gradient_angle = 55.0
    except Exception: pass
    return s

def _spc(run, spc):
    if spc: run._r.get_or_add_rPr().set('spc', str(int(spc)))

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf, m, 0)
    return tb, tf

def para(p, runs, align=PP_ALIGN.LEFT, sa=0, sb=0, line=None):
    p.alignment = align
    if sa is not None: p.space_after = Pt(sa)
    if sb is not None: p.space_before = Pt(sb)
    if line is not None: p.line_spacing = line
    for text, size, color, font, bold, italic, spc in runs:
        r = p.add_run(); r.text = text; f = r.font
        f.size = Pt(size); f.name = font; f.bold = bold; f.italic = italic
        f.color.rgb = color; _spc(r, spc)
    return p

def label(slide, x, y, text, color=BRONZE):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Pt(5), Inches(0.32), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background(); ln.shadow.inherit = False
    tb, tf = textbox(slide, x + Inches(0.46), y, Inches(10.5), Inches(0.35))
    para(tf.paragraphs[0], [(text, 11.5, color, SANS, True, False, 320)])

def monogram(slide, ox, oy, h_in, color):
    s = Inches(h_in) / 50.0
    def rect(x, w):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(ox+x*s), int(oy), int(w*s), int(50*s))
        sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False
    rect(0,10); rect(37.4,10)
    pts = [(25.4,0),(35.4,0),(22,50),(12,50)]
    fb = slide.shapes.build_freeform(int(ox+pts[0][0]*s), int(oy+pts[0][1]*s), scale=1.0)
    fb.add_line_segments([(int(ox+px*s),int(oy+py*s)) for px,py in pts[1:]], close=True)
    sh = fb.convert_to_shape(); sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False

def footer(slide, dark=True, page=None):
    color = IVORY_DIM if dark else BODY_GRAY
    mono = IVORY if dark else INK_TEXT
    y = SH - Inches(0.62)
    monogram(slide, Inches(0.7), y, 0.20, mono)
    tb, tf = textbox(slide, Inches(0.95), y - Pt(2), Inches(6), Inches(0.3))
    para(tf.paragraphs[0], [("MANDOR  ·  ASSIS GESTÃO", 10, color, SANS, True, False, 280)])
    if page:
        tb2, tf2 = textbox(slide, SW - Inches(2.2), y - Pt(2), Inches(1.5), Inches(0.3))
        para(tf2.paragraphs[0], [(f"{page:02d} / {TOTAL:02d}", 10, color, SANS, False, False, 200)], align=PP_ALIGN.RIGHT)

def notes(slide, ancoras, fazer=None, naofazer=None):
    tf = slide.notes_slide.notes_text_frame
    tf.text = "ANCORAS (fale com suas palavras, nao decore)"
    for a in ancoras:
        p = tf.add_paragraph(); p.text = "• " + a
    if fazer:
        p = tf.add_paragraph(); p.text = ""
        p = tf.add_paragraph(); p.text = "FAZER"; p.font.bold = True
        for a in fazer:
            p = tf.add_paragraph(); p.text = "+ " + a
    if naofazer:
        p = tf.add_paragraph(); p.text = ""
        p = tf.add_paragraph(); p.text = "NAO FAZER"; p.font.bold = True
        for a in naofazer:
            p = tf.add_paragraph(); p.text = "x " + a

MX = Inches(0.95)

# ============ SLIDE 1 - CAPA ============
s = add_slide(dark=True)
monogram(s, MX, Inches(1.45), 0.66, IVORY)
tb, tf = textbox(s, MX + Inches(0.82), Inches(1.5), Inches(7), Inches(0.7))
para(tf.paragraphs[0], [("MANDOR", 32, IVORY, SERIF, False, False, 620)])
label(s, MX, Inches(2.95), "CONVERSA INSTITUCIONAL  ·  ASSIS GESTÃO", BRONZE_LT)
tb, tf = textbox(s, MX, Inches(3.5), Inches(11), Inches(2.0))
para(tf.paragraphs[0],
     [("Infraestrutura analítica institucional para ", 40, IVORY, SERIF, False, False, 0),
      ("capital privado", 40, BRONZE_LT, SERIF, False, True, 0)], line=1.05)
tb, tf = textbox(s, MX, Inches(5.35), Inches(10.5), Inches(0.8))
para(tf.paragraphs[0], [("Onde o Mandor amplifica a operação da Assis, e onde não. Uma conversa antes de qualquer demonstração.",
                         16, IVORY_DIM, SERIF, False, False, 0)], line=1.25)
footer(s, dark=True, page=1)
notes(s,
 ["Abrir SEM slide na cabeça: o laptop já está logado no painel, em segundo plano.",
  "Reconhecer a autoridade dos dois sem bajular: Alexandre, 22 anos estruturando para empresário premium; Jamerson, a frente jurídica e de palco.",
  "Sinalizar que o tempo é deles: 30 minutos, você vai direto.",
  "Pedir permissão para começar pelo diagnóstico, não pela ferramenta. Postura de par.",
  "Frase-guia (com suas palavras): 'em vez de apresentar, queria começar entendendo onde a Assis perde mais tempo hoje, e mostrar onde o Mandor encaixa.'"],
 fazer=["Falar devagar. Deixar silêncio depois das perguntas.",
        "Deixar claro que você não veio vender, veio propor rodar junto."],
 naofazer=["Abrir com a história de como o Mandor nasceu.",
           "Dizer 'obrigado pela oportunidade' (registro de vendedor pedindo audiência).",
           "Chamar o Mandor de plataforma, ferramenta ou software. É rede cognitiva."])

# ============ SLIDE 2 - FIT HONESTO 5 de 12 ============
s = add_slide(dark=False)
label(s, MX, Inches(0.85), "ONDE O MANDOR ENCAIXA NA ASSIS", BRONZE)
tb, tf = textbox(s, MX, Inches(1.4), Inches(11.4), Inches(1.0))
para(tf.paragraphs[0],
     [("Dos 12 serviços da Assis, o Mandor amplifica ", 31, INK_TEXT, SERIF, False, False, 0),
      ("cinco", 31, BRONZE, SERIF, False, True, 0),
      (". Os outros sete seguem com vocês.", 31, INK_TEXT, SERIF, False, False, 0)], line=1.08)
# coluna esquerda: 5 amplifica
tb, tf = textbox(s, MX, Inches(2.85), Inches(6.4), Inches(0.4))
para(tf.paragraphs[0], [("O MANDOR AMPLIFICA  ·  5 DE 12", 12, BRONZE, SANS, True, False, 200)])
cinco = ["M&A (Fusões e Aquisições)", "Due Diligence", "Valuation",
         "Estruturação de Capital", "Diagnóstico Empresarial (parte analítica)"]
yy = Inches(3.45)
for i, item in enumerate(cinco):
    tk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, yy + Inches(0.56)*i + Pt(4), Inches(0.2), Pt(2.4))
    tk.fill.solid(); tk.fill.fore_color.rgb = BRONZE; tk.line.fill.background(); tk.shadow.inherit = False
    tb, tf = textbox(s, MX + Inches(0.38), yy + Inches(0.56)*i, Inches(6.0), Inches(0.5))
    para(tf.paragraphs[0], [(item, 16.5, INK_TEXT, SERIF, False, False, 0)], line=1.0)
# divisor
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.65), Inches(2.95), Pt(1.2), Inches(3.4))
dv.fill.solid(); dv.fill.fore_color.rgb = RGBColor(0xCF,0xC7,0xBB); dv.line.fill.background(); dv.shadow.inherit = False
# coluna direita: 7 seguem
tb, tf = textbox(s, Inches(8.05), Inches(2.85), Inches(4.6), Inches(0.4))
para(tf.paragraphs[0], [("SEGUE 100% COM A ASSIS  ·  7 DE 12", 12, BODY_GRAY, SANS, True, False, 200)])
sete = "Governança · Compliance contínuo · Auditoria contábil · Administração Judicial · Treinamento de Gestão · Calendário Contábil · Balanced Scorecard"
tb, tf = textbox(s, Inches(8.05), Inches(3.45), Inches(4.6), Inches(2.6))
para(tf.paragraphs[0], [(sete, 15, BODY_GRAY, SERIF, False, False, 0)], line=1.5)
footer(s, dark=False, page=2)
notes(s,
 ["Este é o slide mais forte da reunião: prova que você estudou a Assis, e que não promete demais.",
  "Dizer abertamente: 'sete dos doze serviços de vocês o Mandor não toca. Continuam sendo vocês.'",
  "Os cinco que batem são justamente os mais transacionais e de maior margem da Assis.",
  "O enquadramento: o Mandor não substitui a Assis, amplifica a capacidade nos cinco e libera a equipe para os outros sete.",
  "Inverte o jogo: você não pede para vender, você identificou onde dá para trabalhar junto."],
 fazer=["Se ele perguntar 'e a profundidade dos cinco?', responder: a diferença não é profundidade, é velocidade com a profundidade preservada. A autoria continua sendo da Assis.",
        "Conectar ao caso real deles (ex.: Avantgarde) se você souber quais dos cinco a operação usou."],
 naofazer=["Dizer que o Mandor cobre os 12. Ele testa item por item e você perde o frame.",
           "Entrar em detalhe de cada um dos sete. Só nomeie e siga."])

# ============ SLIDE 3 - AS 3 PERGUNTAS ============
s = add_slide(dark=True)
label(s, MX, Inches(0.85), "DENTRO DESSES CINCO, TRÊS PERGUNTAS", BRONZE_LT)
tb, tf = textbox(s, MX, Inches(1.45), Inches(11), Inches(1.0))
para(tf.paragraphs[0],
     [("As três perguntas que aparecem em ", 32, IVORY, SERIF, False, False, 0),
      ("toda casa que cresce", 32, BRONZE_LT, SERIF, False, True, 0)], line=1.06)
perg = [("01", "A inteligência da casa fica, ou vai embora com o sócio?"),
        ("02", "O deal chega ao comprador certo, ou depende da agenda de alguém?"),
        ("03", "A capacidade de análise acompanha o volume de oportunidades?")]
yy = Inches(3.1)
for i,(n,q) in enumerate(perg):
    tb, tf = textbox(s, MX, yy + Inches(1.0)*i, Inches(11.4), Inches(0.9))
    para(tf.paragraphs[0], [(n, 22, BRONZE_LT, SERIF, False, False, 0),
                            ("    " + q, 24, IVORY, SERIF, False, False, 0)], line=1.05)
footer(s, dark=True, page=3)
notes(s,
 ["Apresentar as três como perguntas honestas de quem cresce, não como falhas. Não usar dor, problema, gargalo.",
  "01 memória institucional: em 22 anos, quanto do método vive na cabeça da equipe sênior e quanto fica na casa?",
  "02 originação: quantos deals dependeram do sócio bater no telefone para achar o comprador certo?",
  "03 escala: se desse para dar parecer rápido com o rigor da própria mesa, quantos deals a mais entrariam no funil?",
  "Cada pergunta é de validação: deixe ELES responderem. Quanto mais eles falam, melhor.",
  "Transição para a demo: 'os três pontos que vocês descreveram são os que o Mandor foi construído para responder. Posso mostrar, sem deck?'"],
 fazer=["Se ele recusar uma das perguntas ('isso está resolvido aqui'), aceitar e seguir para a próxima. Não insistir."],
 naofazer=["Pitchar a solução em cima da pergunta. Aqui só se abre o tema."])

# ============ SLIDE 4 - DEMO DIVIDER ============
s = add_slide(dark=True)
label(s, MX, Inches(0.9), "DEMONSTRAÇÃO AO VIVO", BRONZE_LT)
tb, tf = textbox(s, MX, Inches(2.6), Inches(11.3), Inches(1.6), anchor=MSO_ANCHOR.MIDDLE)
para(tf.paragraphs[0],
     [("Em vez de explicar, deixa eu ", 38, IVORY, SERIF, False, False, 0),
      ("mostrar", 38, BRONZE_LT, SERIF, False, True, 0),
      (".", 38, IVORY, SERIF, False, False, 0)], line=1.05)
tb, tf = textbox(s, MX, Inches(4.4), Inches(11), Inches(1.0))
para(tf.paragraphs[0], [("Lumina Diagnósticos. Venda de controle, medicina diagnóstica no Sul. Um caso do mundo de vocês, rodado por dentro do Mandor.",
                         16.5, IVORY_DIM, SERIF, False, False, 0)], line=1.3)
footer(s, dark=True, page=4)
notes(s,
 ["AQUI VOCÊ SAI DOS SLIDES E VAI PARA O PAINEL. Veja a folha 'Mapa da Reunião' para o caminho exato da demo.",
  "Demo ancorada nos 3 gargalos, na mesma ordem (cerca de 10 minutos, sem voltar ao deck).",
  "Gargalo 01 -> aba Aprendizados do escritório + Benchmarks (o método e o critério ficam na casa).",
  "Gargalo 02 -> Invest Match do Lumina (LEIA OS NOMES NA TELA, não decore: são investidores de saúde/PE, com score e justificativa).",
  "Gargalo 03 -> Dossiê do Lumina: DRS 83/100, valuation, riscos, comprador mapeado, teaser. Cada conclusão com fonte e página (lineage).",
  "O dossiê é o momento de decisão: navegue devagar para ele ler trechos."],
 fazer=["Antes da reunião: abrir o Lumina no painel e confirmar que Invest Match, Benchmarks e Dossiê carregam. A demo foi re-semeada hoje."],
 naofazer=["Recitar 'Kinea 92, Pátria 91'. Esses eram de outro deal (crédito). No Lumina os nomes são outros: leia a tela."])

# ============ SLIDE 5 - ESTADO RESOLVIDO ============
s = add_slide(dark=False)
label(s, MX, Inches(0.85), "O QUE MUDA QUANDO ISSO ESTÁ RESOLVIDO", BRONZE)
tb, tf = textbox(s, MX, Inches(1.45), Inches(11), Inches(0.9))
para(tf.paragraphs[0], [("Para uma casa como a Assis, na prática:", 30, INK_TEXT, SERIF, False, False, 0)], line=1.05)
res = ["A inteligência analítica da casa fica na casa, mesmo quando um sócio sai.",
       "Cada deal sai do diagnóstico já sabendo para quem levar, com score e justificativa.",
       "O mesmo critério da mesa de vocês vale no primeiro e no quadragésimo deal do mês."]
yy = Inches(2.95)
for i,o in enumerate(res):
    tk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, yy + Inches(0.85)*i + Pt(6), Inches(0.22), Pt(2.6))
    tk.fill.solid(); tk.fill.fore_color.rgb = BRONZE; tk.line.fill.background(); tk.shadow.inherit = False
    tb, tf = textbox(s, MX + Inches(0.42), yy + Inches(0.85)*i, Inches(10.6), Inches(0.8))
    para(tf.paragraphs[0], [(o, 19, INK_TEXT, SERIF, False, False, 0)], line=1.1)
footer(s, dark=False, page=5)
notes(s,
 ["Mudar o tom: do diagnóstico para a imagem do estado desejado. Descrever, não prometer.",
  "Deixe que eles se projetem nessa imagem depois de terem visto a demo.",
  "Amarra a autoria: a velocidade é do Mandor, o critério e a assinatura continuam sendo da Assis.",
  "É a ponte natural para o próximo passo."],
 naofazer=["Prometer resultado de receita ou número. Aqui é estado, não promessa."])

# ============ SLIDE 6 - PROXIMO PASSO (PILOTO) ============
s = add_slide(dark=False)
label(s, MX, Inches(0.85), "PRÓXIMO PASSO", BRONZE)
tb, tf = textbox(s, MX, Inches(1.45), Inches(11), Inches(0.9))
para(tf.paragraphs[0],
     [("Um piloto sobre um caso ", 31, INK_TEXT, SERIF, False, False, 0),
      ("que vocês conhecem", 31, BRONZE, SERIF, False, True, 0),
      (".", 31, INK_TEXT, SERIF, False, False, 0)], line=1.05)
passos = [("01", "Vocês me passam um caso anonimizado que a Assis estruturou nos últimos meses."),
          ("02", "Eu rodo dentro do Mandor e preparo o dossiê sobre material que vocês conhecem."),
          ("03", "Call em cerca de dez dias para verem funcionando. O desenho comercial conversamos a partir daí.")]
yy = Inches(2.95)
for i,(n,t) in enumerate(passos):
    tb, tf = textbox(s, MX, yy + Inches(0.92)*i, Inches(11.4), Inches(0.85))
    para(tf.paragraphs[0], [(n, 19, BRONZE, SERIF, False, False, 0),
                            ("    " + t, 18, INK_TEXT, SERIF, False, False, 0)], line=1.1)
footer(s, dark=False, page=6)
notes(s,
 ["O pedido é só o piloto. NADA de número (preço ou comissão) nesta reunião.",
  "Se perguntarem de valores: 'o desenho comercial a gente fecha depois que vocês virem rodando sobre um caso de vocês'.",
  "Pedir um caso anonimizado, de preferência do tier premium (ex.: algo no perfil Avantgarde).",
  "Combinar a call de dez dias ali, com data tentativa.",
  "Se ele pedir tempo, aceitar e propor mandar o NDA já, para destravar quando decidirem."],
 fazer=["Sair com: caso anonimizado prometido, call agendada, NDA combinado."],
 naofazer=["Forçar fechamento. Sênior que sente pressão fecha a porta. Relação intacta vale mais que um sim hoje."])

# ============ SLIDE 7 - FECHAMENTO ============
s = add_slide(dark=True)
monogram(s, MX, Inches(2.5), 0.6, IVORY)
tb, tf = textbox(s, MX + Inches(0.78), Inches(2.55), Inches(6), Inches(0.6))
para(tf.paragraphs[0], [("MANDOR", 28, IVORY, SERIF, False, False, 620)])
tb, tf = textbox(s, MX, Inches(3.6), Inches(11), Inches(1.2))
para(tf.paragraphs[0],
     [("Não vim vender. Vim propor que a gente ", 30, IVORY, SERIF, False, False, 0),
      ("rode junto", 30, BRONZE_LT, SERIF, False, True, 0),
      (".", 30, IVORY, SERIF, False, False, 0)], line=1.1)
tb, tf = textbox(s, MX, Inches(5.0), Inches(8), Inches(0.5))
para(tf.paragraphs[0], [("www.mandor.com.br", 16, BRONZE_LT, SANS, True, False, 120)])
footer(s, dark=True, page=7)
notes(s,
 ["Fechar devolvendo a palavra a eles. Postura de par até o fim.",
  "Se a relação ficou boa mas sem sim, está ótimo: o objetivo é virar aliado institucional, não fechar hoje.",
  "Cartão no bolso, NDA pronto para enviar no mesmo dia se aceitarem o piloto."])

out = "/Users/renan/Desktop/rr7x-portal/Mandor-Apresentacao/Mandor-Assis-MA.pptx"
prs.save(out)
print("OK ->", out)
